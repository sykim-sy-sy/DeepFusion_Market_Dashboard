from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ppt(output_path):
    prs = Presentation()

    # Define some colors (DFAI Blue)
    dfai_blue = RGBColor(0, 102, 204)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "홈페이지 개편 추진 현황 및 개편안"
    subtitle.text = "Website Reorganization Progress & Plan\n2026. 03. 12."

    # 2. Current Status & Feedback
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. 현재 진행 상황 및 피드백"
    body = slide.placeholders[1].text_frame
    body.text = "추진 현황"
    body.add_paragraph().text = "• 1월 5일: 1차 디자인 시안 검토 및 피드백 전달"
    body.add_paragraph().text = "• 주요 피드백:"
    p = body.add_paragraph()
    p.text = "    - 파란 계열(CI Identity) 부족: 회사 정체성 미비"
    p = body.add_paragraph()
    p.text = "    - 핵심 기술 및 슬로건 노출 약화"
    body.add_paragraph().text = "• 현재 단계: 업체(제작사)로부터 핵심 문구 및 기술 표현 이미지 설명 요청 수신"

    # 3. 개편 배경 및 필요성 (From PDF)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. 개편 배경 및 필요성"
    body = slide.placeholders[1].text_frame
    body.text = "기존 홈페이지의 한계"
    body.add_paragraph().text = "• 정보 구조 분산: B2B 고객의 핵심 제품/기술 이해 난해"
    body.add_paragraph().text = "• 영업 리드 확보 어려움: 문의(Contact) 체계 부재"
    body.add_paragraph().text = "• 글로벌 대응 미흡: 모바일 및 다국어 지원 한계"

    # 4. 개편 목표 (From PDF)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. 개편 목표 및 방향"
    body = slide.placeholders[1].text_frame
    body.text = "브랜드 정체성 강화"
    body.add_paragraph().text = "• Perceptive Sensor Fusion System Company 정체성 명확화"
    body.add_paragraph().text = "• 영업·마케팅 플랫폼으로의 전환: 단순 소개가 아닌 '플랫폼' 역할"
    body.add_paragraph().text = "• 신뢰성 확보: 핵심 기술(RAPA-R) 및 파트너 레퍼런스 가시화"

    # 5. 메뉴 구조 기획 (From PDF)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4. 홈페이지 메뉴 구조 (안)"
    body = slide.placeholders[1].text_frame
    body.text = "상단 메뉴 5개 이내 단순화"
    body.add_paragraph().text = "1. About: 회사 소개, 파트너, 테스티모니얼"
    body.add_paragraph().text = "2. Solutions: RAPA-R, Multi Radar SLAM"
    body.add_paragraph().text = "3. Technology: Radar AI 기술, 산업군(Industries)"
    body.add_paragraph().text = "4. Newsroom: IR/PR, Media, Awards, Blog 통합"
    body.add_paragraph().text = "5. Careers: 채용 정보 및 기업 문화"

    # 6. 업체 요청 대응: 기술 이미지 가이드 (Appendix)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "5. 기술/슬로건 이미지 표현 가이드 (요청 대응)"
    body = slide.placeholders[1].text_frame
    body.text = "심상(Visual Concept)"
    body.add_paragraph().text = "• 핵심 컬러: DFAI 파란계열(Primary) + 다크 모드(Professional UX)"
    body.add_paragraph().text = "• 이미지 1 (Perceptive Fusion): 안개/폭우 속의 도로를 4D 포인트 클라우드로 뚫고 지나가는 역동적인 뷰"
    body.add_paragraph().text = "• 이미지 2 (RAPA-R): 레이더 신호가 필라(Pillar) 구조로 변환되어 실시간으로 객체를 인지하는 AI 아키텍처 비주얼"
    body.add_paragraph().text = "• 이미지 3 (Identity): Perceptive Sensor Fusion 문구와 함께 레이더+카메라 데이터가 융합되는 큐브 형태의 비주얼"

    prs.save(output_path)
    print(f"PPT generated: {output_path}")

if __name__ == "__main__":
    create_ppt(r"C:\Users\KIM SIYE\OneDrive - 딥퓨전에이아이주식회사\바탕 화면\SIYE\Marketing_Team\홈페이지_개편_추진_현황_및_기획안.pptx")
